#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
MizanQuant — Investor Briefing PPTX Exporter
============================================

Builds a native, fully-editable 6-slide .pptx that mirrors the HTML deck at
slides/index.html. Every block is a real PowerPoint shape (rectangle, oval,
line, text box) — not an image — so an analyst can swap copy, retype numbers,
or restyle a single tile without breaking the layout.

Slide canvas: 16:9 at 13.333" × 7.5" (1920×1080 px @ 144 dpi).

Usage
-----
    pip install python-pptx
    python export-deck.py [--out mizanquant-briefing.pptx]

Output
------
A single .pptx written to the same directory (or `--out`).

Design tokens
-------------
All colors, fonts, and spacing constants are pulled from
../colors_and_type.css and live in the BRAND block below. Change a token there
and the entire deck restyles.

Requirements
------------
- python-pptx >= 0.6.21
- DM Sans, JetBrains Mono, and Cairo installed on the OS that opens the file
  (PowerPoint falls back gracefully; for embedded fonts use the GUI's File →
  Options → Save → Embed fonts).
"""

from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# =============================================================================
# BRAND — single source of truth, mirrors colors_and_type.css
# =============================================================================

# Slide canvas (16:9, 1920×1080 @ 144dpi)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — warm near-black + warm gold
BG_ROOT      = RGBColor(0x0C, 0x0C, 0x10)
BG_SURFACE   = RGBColor(0x14, 0x14, 0x1A)
BG_RAISED    = RGBColor(0x1C, 0x1C, 0x26)
BG_OVERLAY   = RGBColor(0x24, 0x24, 0x30)

TEXT_PRIMARY   = RGBColor(0xE4, 0xE4, 0xEC)
TEXT_SECONDARY = RGBColor(0x94, 0x94, 0xA0)
TEXT_MUTED     = RGBColor(0x5C, 0x5C, 0x6A)
TEXT_DISABLED  = RGBColor(0x3C, 0x3C, 0x48)

ACCENT       = RGBColor(0xC8, 0x96, 0x3E)   # the gold
ACCENT_HOVER = RGBColor(0xD4, 0xA8, 0x53)
ACCENT_DIM   = RGBColor(0x2A, 0x22, 0x16)   # ~10% gold over bg-root, opaque approximation

POSITIVE     = RGBColor(0x4A, 0xDE, 0x80)
NEGATIVE     = RGBColor(0xF8, 0x71, 0x71)
WARNING      = RGBColor(0xFB, 0xBF, 0x24)

BORDER_DEFAULT = RGBColor(0x2A, 0x2A, 0x32)   # ~6% white over bg-root
BORDER_SUBTLE  = RGBColor(0x1A, 0x1A, 0x20)   # ~4% white over bg-root

# Type stacks
FONT_DISPLAY = "DM Sans"
FONT_BODY    = "DM Sans"      # PPT does not use system stacks; pick a single face
FONT_MONO    = "JetBrains Mono"
FONT_ARABIC  = "Cairo"

# Layout pads (mirror landing/deck CSS: 96px @ 144dpi ≈ 0.667 in)
PAD_X = Inches(0.667)
PAD_Y = Inches(0.5)


# =============================================================================
# UTILITY — shape & text primitives
# =============================================================================

def fill_solid(shape, rgb: RGBColor) -> None:
    """Set a shape's fill to a solid color."""
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def fill_none(shape) -> None:
    shape.fill.background()


def line_solid(shape, rgb: RGBColor, width_pt: float = 0.75) -> None:
    """Set a shape's outline to a solid color and width."""
    line = shape.line
    line.color.rgb = rgb
    line.width = Pt(width_pt)


def line_none(shape) -> None:
    shape.line.fill.background()


def add_rect(
    slide,
    left, top, width, height,
    *,
    fill: Optional[RGBColor] = BG_SURFACE,
    border: Optional[RGBColor] = None,
    border_w_pt: float = 0.75,
    radius: float = 0.04,
):
    """Add a (optionally rounded) rectangle shape."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius > 0 else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        # adjustment 0 controls corner radius (0 = sharp, 0.5 = max)
        shape.adjustments[0] = radius

    if fill is None:
        fill_none(shape)
    else:
        fill_solid(shape, fill)

    if border is None:
        line_none(shape)
    else:
        line_solid(shape, border, border_w_pt)

    shape.shadow.inherit = False
    return shape


def add_oval(slide, left, top, width, height, *, fill=ACCENT_DIM, border=ACCENT, border_w_pt=0.75):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if fill is None:
        fill_none(shape)
    else:
        fill_solid(shape, fill)
    if border is None:
        line_none(shape)
    else:
        line_solid(shape, border, border_w_pt)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, *, color: RGBColor = BORDER_DEFAULT, width_pt: float = 0.5):
    shape = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line_solid(shape, color, width_pt)
    return shape


@dataclass
class TextStyle:
    font: str = FONT_BODY
    size_pt: float = 12
    bold: bool = False
    italic: bool = False
    color: RGBColor = TEXT_PRIMARY
    align: PP_ALIGN = PP_ALIGN.LEFT
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP
    tracking_pt: float = 0  # letter spacing in points
    line_spacing: float = 1.15


def add_text(
    slide,
    left, top, width, height,
    text: str,
    style: TextStyle,
    *,
    bg: Optional[RGBColor] = None,
):
    """Add a single-style text box. Returns the shape so callers can post-modify."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = style.anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)

    if bg is not None:
        fill_solid(box, bg)
    else:
        fill_none(box)
    line_none(box)
    box.shadow.inherit = False

    p = tf.paragraphs[0]
    p.alignment = style.align
    p.line_spacing = style.line_spacing
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = style.font
    f.size = Pt(style.size_pt)
    f.bold = style.bold
    f.italic = style.italic
    f.color.rgb = style.color
    return box


def add_rich_text(
    slide,
    left, top, width, height,
    runs: Sequence[tuple],   # iterable of (text, TextStyle) — one paragraph
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.15,
    bg: Optional[RGBColor] = None,
):
    """Add a text box whose single paragraph mixes multiple styles per run."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    if bg is not None:
        fill_solid(box, bg)
    else:
        fill_none(box)
    line_none(box)
    box.shadow.inherit = False
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for text, style in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = style.font
        f.size = Pt(style.size_pt)
        f.bold = style.bold
        f.italic = style.italic
        f.color.rgb = style.color
    return box


# Common styles ---------------------------------------------------------------
EYEBROW_STYLE = TextStyle(font=FONT_BODY, size_pt=11, bold=True, color=ACCENT)
TITLE_STYLE   = TextStyle(font=FONT_DISPLAY, size_pt=44, bold=True, color=TEXT_PRIMARY, line_spacing=1.05)
ACCENT_TITLE  = TextStyle(font=FONT_DISPLAY, size_pt=44, bold=True, color=ACCENT,        line_spacing=1.05)
SUBTITLE_STYLE = TextStyle(font=FONT_DISPLAY, size_pt=18, color=TEXT_SECONDARY, line_spacing=1.35)
LABEL_STYLE    = TextStyle(font=FONT_BODY, size_pt=8.5, bold=True, color=TEXT_MUTED)
METRIC_STYLE   = TextStyle(font=FONT_MONO, size_pt=34, bold=True, color=TEXT_PRIMARY, line_spacing=1.0)
BODY_STYLE     = TextStyle(font=FONT_BODY, size_pt=11, color=TEXT_SECONDARY, line_spacing=1.5)
MONO_SMALL     = TextStyle(font=FONT_MONO, size_pt=10, color=TEXT_MUTED)
ARABIC_STYLE   = TextStyle(font=FONT_ARABIC, size_pt=18, bold=True, color=ACCENT)


# =============================================================================
# REUSABLE COMPOSITES — brand mark, slide chrome, eyebrow row
# =============================================================================

def draw_brand_mark(slide, left, top, size=Inches(0.5), color=ACCENT, halo: bool = True):
    """Mīm-trace icon: ring + ascending dot-to-dot."""
    if halo:
        add_rect(slide, left, top, size, size, fill=ACCENT_DIM, border=ACCENT, border_w_pt=0.5, radius=0.15)

    # Coordinates inside the tile (64×64 SVG viewBox)
    inner = size * 0.85
    ox = left + (size - inner) / 2
    oy = top + (size - inner) / 2
    s = inner / Emu(914400) * Emu(914400)  # keep type

    def scale(x):
        return ox + Emu(int(x / 64 * inner)) if False else ox + (x / 64) * inner

    def scale_y(y):
        return oy + (y / 64) * inner

    # ring
    ring_d = (9.5 * 2 / 64) * inner
    ring = add_oval(
        slide,
        ox + (18 - 9.5) / 64 * inner, oy + (34 - 9.5) / 64 * inner,
        ring_d, ring_d,
        fill=None, border=color, border_w_pt=2.2,
    )
    # chart segments
    segs = [(27, 40, 36, 30), (36, 30, 44, 36), (44, 36, 54, 22)]
    for x1, y1, x2, y2 in segs:
        add_line(
            slide,
            ox + x1 / 64 * inner, oy + y1 / 64 * inner,
            ox + x2 / 64 * inner, oy + y2 / 64 * inner,
            color=color, width_pt=2.2,
        )
    # terminal dot
    dot_d = (3 * 2 / 64) * inner
    add_oval(
        slide,
        ox + (54 - 3) / 64 * inner, oy + (22 - 3) / 64 * inner,
        dot_d, dot_d,
        fill=color, border=None,
    )


def draw_eyebrow(slide, left, top, text: str):
    """Eyebrow row: gold dot + uppercase eyebrow text."""
    dot_size = Inches(0.08)
    add_oval(slide, left, top + Inches(0.025), dot_size, dot_size, fill=ACCENT, border=None)
    add_text(
        slide,
        left + Inches(0.16), top - Inches(0.02), Inches(8), Inches(0.25),
        text.upper(),
        TextStyle(font=FONT_BODY, size_pt=11, bold=True, color=ACCENT),
    )


def draw_slide_chrome(slide, page_no: int, total: int = 6):
    """Bottom bar: small wordmark + page number."""
    y = SLIDE_H - Inches(0.45)
    add_rich_text(
        slide, PAD_X, y, Inches(6), Inches(0.25),
        runs=[
            ("mizan", TextStyle(font=FONT_DISPLAY, size_pt=10, bold=True, color=TEXT_SECONDARY)),
            ("quant", TextStyle(font=FONT_DISPLAY, size_pt=10, bold=True, color=ACCENT)),
            ("  ·  briefing", TextStyle(font=FONT_DISPLAY, size_pt=10, color=TEXT_MUTED)),
        ],
    )
    add_text(
        slide,
        SLIDE_W - PAD_X - Inches(2), y, Inches(2), Inches(0.25),
        f"{page_no:02d} / {total:02d}",
        TextStyle(font=FONT_MONO, size_pt=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT),
    )


def fill_slide_bg(slide, color=BG_ROOT):
    """Set the slide background to a solid color (creates a full-bleed rect)."""
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill=color, border=None, radius=0)


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_title_slide(prs: Presentation) -> None:
    """Slide 01 — Title page."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    fill_slide_bg(slide)

    # Faint gold glow in top-right (approximated as a layered rounded rect)
    add_rect(
        slide,
        SLIDE_W - Inches(6), Inches(-2), Inches(7), Inches(5),
        fill=RGBColor(0x18, 0x14, 0x18), border=None, radius=0.5,
    )

    # ── top brand row ──
    brand_y = Inches(0.6)
    draw_brand_mark(slide, PAD_X, brand_y, size=Inches(0.55))
    add_rich_text(
        slide, PAD_X + Inches(0.75), brand_y - Inches(0.04), Inches(4.5), Inches(0.4),
        runs=[
            ("mizan", TextStyle(font=FONT_DISPLAY, size_pt=22, bold=True, color=TEXT_PRIMARY)),
            ("quant", TextStyle(font=FONT_DISPLAY, size_pt=22, bold=True, color=ACCENT)),
        ],
    )
    add_text(
        slide, PAD_X + Inches(0.75), brand_y + Inches(0.30), Inches(4.5), Inches(0.3),
        "ميزان كوانت",
        TextStyle(font=FONT_ARABIC, size_pt=14, bold=True, color=TEXT_SECONDARY),
    )

    # ── top-right meta ──
    add_rich_text(
        slide, SLIDE_W - PAD_X - Inches(4), brand_y, Inches(4), Inches(0.5),
        runs=[("Confidential · Investor briefing\nQ2 2026  ·  NDA on file",
               TextStyle(font=FONT_MONO, size_pt=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT))],
        align=PP_ALIGN.RIGHT,
    )

    # ── HEADLINE block ──
    head_top = Inches(2.4)
    add_rich_text(
        slide, PAD_X, head_top, SLIDE_W - PAD_X * 2, Inches(2.4),
        runs=[
            ("The convergence of ",       TextStyle(font=FONT_DISPLAY, size_pt=56, bold=True, color=TEXT_PRIMARY, line_spacing=1.04)),
            ("Sharia governance",         TextStyle(font=FONT_DISPLAY, size_pt=56, bold=True, color=ACCENT,        line_spacing=1.04)),
            ("\nand ",                    TextStyle(font=FONT_DISPLAY, size_pt=56, bold=True, color=TEXT_PRIMARY, line_spacing=1.04)),
            ("algorithmic alpha.",        TextStyle(font=FONT_DISPLAY, size_pt=56, bold=True, color=ACCENT,        line_spacing=1.04)),
        ],
        line_spacing=1.04,
    )

    add_text(
        slide, PAD_X, head_top + Inches(2.5), SLIDE_W - PAD_X * 2, Inches(0.5),
        "حيث تلتقي حوكمة الشريعة بمنطق الخوارزميات.",
        TextStyle(font=FONT_ARABIC, size_pt=22, bold=True, color=ACCENT),
    )

    # ── bottom strip ──
    bot = SLIDE_H - Inches(1.1)
    add_text(slide, PAD_X, bot, Inches(6), Inches(0.25),
             "BRIEFING AUDIENCE",
             TextStyle(font=FONT_BODY, size_pt=9, bold=True, color=TEXT_MUTED))
    add_text(slide, PAD_X, bot + Inches(0.28), Inches(8), Inches(0.4),
             "Sovereign wealth funds · Sharia-compliant asset managers · Family offices",
             TextStyle(font=FONT_DISPLAY, size_pt=14, bold=True, color=TEXT_PRIMARY))

    add_rich_text(
        slide, SLIDE_W - PAD_X - Inches(4), bot, Inches(4), Inches(0.7),
        runs=[("Faisal — MizanQuant\ninstitutional@mizanquant.com",
               TextStyle(font=FONT_MONO, size_pt=9, color=TEXT_MUTED, align=PP_ALIGN.RIGHT))],
        align=PP_ALIGN.RIGHT,
    )


def build_architecture_slide(prs: Presentation) -> None:
    """Slide 02 — 5-stage architecture loop."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)

    # Header
    draw_eyebrow(slide, PAD_X, Inches(0.55), "Architecture")
    add_rich_text(
        slide, PAD_X, Inches(0.95), Inches(11.5), Inches(1.4),
        runs=[
            ("A five-layer loop, supervised by a ", TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=TEXT_PRIMARY, line_spacing=1.08)),
            ("standing risk desk.",                  TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=ACCENT,        line_spacing=1.08)),
        ],
        line_spacing=1.08,
    )
    add_text(
        slide, PAD_X, Inches(2.15), Inches(11.5), Inches(0.7),
        "Every trade begins as a halal-screened candidate and ends as a stress-tested position with a Kelly-sized risk budget. The loop runs end-to-end six times a day.",
        TextStyle(font=FONT_DISPLAY, size_pt=14, color=TEXT_SECONDARY, line_spacing=1.45),
    )

    # 5-stage grid
    stages = [
        ("01 / GATE",      "Halal screener",        "AAOIFI four-screen gate.\nDebt · interest · cash ·\nharam industries.\nRescanned daily."),
        ("02 / SIGNAL",    "Forecast lab",          "14 forecast models vote\nper session — transformer,\nGRU, LSTM, CNN, classical\nensemble."),
        ("03 / CONSENSUS", "AI consensus",          "Confidence-weighted vote.\nSTRONG BUY requires\n≥ 11 of 14 votes plus\nscore ≥ 90."),
        ("04 / SIZE",      "Kelly allocator",       "Continuous Kelly with\nsample-size shrinkage,\nhalf-Kelly default, 25%\nhard ceiling."),
        ("05 / EXECUTE",   "Risk desk · execute",   "Pre-trade VaR check,\nguard pass, kill-switch,\nbroker handoff (Alpaca /\nIBKR), audit log."),
    ]

    grid_top = Inches(3.4)
    grid_h   = Inches(3.0)
    n = 5
    gap = Inches(0.14)
    total_w = SLIDE_W - PAD_X * 2
    col_w = (total_w - gap * (n - 1)) / n

    for i, (no, name, desc) in enumerate(stages):
        x = PAD_X + (col_w + gap) * i
        card = add_rect(slide, x, grid_top, col_w, grid_h, fill=BG_SURFACE, border=BORDER_DEFAULT, border_w_pt=0.5, radius=0.05)
        # accent stripe at top
        add_rect(slide, x, grid_top, col_w, Inches(0.05), fill=ACCENT, border=None, radius=0)

        ix, iy = x + Inches(0.22), grid_top + Inches(0.22)
        add_text(slide, ix, iy, col_w - Inches(0.4), Inches(0.25), no,
                 TextStyle(font=FONT_MONO, size_pt=9, bold=True, color=ACCENT))
        add_text(slide, ix, iy + Inches(0.36), col_w - Inches(0.4), Inches(0.5), name,
                 TextStyle(font=FONT_DISPLAY, size_pt=14, bold=True, color=TEXT_PRIMARY))
        add_text(slide, ix, iy + Inches(0.9), col_w - Inches(0.4), Inches(1.9), desc,
                 TextStyle(font=FONT_BODY, size_pt=10, color=TEXT_SECONDARY, line_spacing=1.4))

    # Loop feedback label
    add_text(
        slide,
        PAD_X, grid_top + grid_h + Inches(0.35),
        SLIDE_W - PAD_X * 2, Inches(0.3),
        "↻   FEEDBACK LOOP · REALIZED PnL · MODEL ACCURACY · REGIME UPDATES · HALAL RE-VERIFY   ↻",
        TextStyle(font=FONT_MONO, size_pt=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER),
    )

    draw_slide_chrome(slide, 2)


def build_halal_slide(prs: Presentation) -> None:
    """Slide 03 — AAOIFI four-screen + NVDA specimen."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)

    # Header
    draw_eyebrow(slide, PAD_X, Inches(0.55), "Sharia compliance · core metric")
    add_rich_text(
        slide, PAD_X, Inches(0.95), Inches(11.5), Inches(1.4),
        runs=[
            ("Compliance is the ",            TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=TEXT_PRIMARY, line_spacing=1.08)),
            ("first gate",                    TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=ACCENT,        line_spacing=1.08)),
            (", not a footnote.",             TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=TEXT_PRIMARY, line_spacing=1.08)),
        ],
        line_spacing=1.08,
    )

    add_text(
        slide, PAD_X, Inches(2.15), Inches(11.5), Inches(0.6),
        "Every symbol passes AAOIFI's four-screen test before it can become a candidate signal. No override path bypasses this gate.",
        TextStyle(font=FONT_DISPLAY, size_pt=14, color=TEXT_SECONDARY, line_spacing=1.45),
    )

    # Two columns: left=screen list, right=NVDA specimen
    col_top = Inches(3.0)
    col_h   = Inches(3.5)
    left_x  = PAD_X
    left_w  = Inches(6.2)
    right_x = PAD_X + left_w + Inches(0.5)
    right_w = SLIDE_W - PAD_X - right_x

    # LEFT — screen list
    screens = [
        (1, "Debt screen",        "Total Debt / Market Cap",       "< 33%"),
        (2, "Interest income",    "Interest Income / Revenue",     "< 5%"),
        (3, "Cash & receivables", "(Cash + AR) / Market Cap",      "< 33%"),
        (4, "Haram industries",   "AAOIFI excluded sector list",   "Excluded"),
    ]
    row_h = Inches(0.78)
    row_gap = Inches(0.08)
    for i, (n, name, formula, threshold) in enumerate(screens):
        y = col_top + (row_h + row_gap) * i
        # row card
        add_rect(slide, left_x, y, left_w, row_h, fill=BG_SURFACE, border=BORDER_DEFAULT, border_w_pt=0.5, radius=0.08)
        # gold left stripe
        add_rect(slide, left_x, y, Inches(0.08), row_h, fill=ACCENT, border=None, radius=0)
        # number bubble
        bn_size = Inches(0.4)
        add_oval(slide, left_x + Inches(0.28), y + (row_h - bn_size) / 2, bn_size, bn_size, fill=ACCENT_DIM, border=ACCENT, border_w_pt=0.5)
        add_text(slide, left_x + Inches(0.28), y + (row_h - bn_size) / 2 - Inches(0.04), bn_size, bn_size + Inches(0.08),
                 str(n),
                 TextStyle(font=FONT_MONO, size_pt=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE))
        # name + formula
        add_text(slide, left_x + Inches(0.85), y + Inches(0.12), left_w - Inches(2.2), Inches(0.3),
                 name,
                 TextStyle(font=FONT_DISPLAY, size_pt=14, bold=True, color=TEXT_PRIMARY))
        add_text(slide, left_x + Inches(0.85), y + Inches(0.40), left_w - Inches(2.2), Inches(0.3),
                 formula,
                 TextStyle(font=FONT_MONO, size_pt=10, color=TEXT_MUTED))
        # threshold
        add_text(slide, left_x + left_w - Inches(1.6), y, Inches(1.45), row_h,
                 threshold,
                 TextStyle(font=FONT_MONO, size_pt=18, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE))

    # RIGHT — NVDA specimen
    spec_h = (row_h + row_gap) * 4 - row_gap
    add_rect(slide, right_x, col_top, right_w, spec_h, fill=BG_SURFACE, border=BORDER_DEFAULT, border_w_pt=0.5, radius=0.05)
    pad = Inches(0.4)
    add_text(slide, right_x + pad, col_top + Inches(0.2), right_w - pad * 2, Inches(0.7),
             "NVDA",
             TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=TEXT_PRIMARY))
    add_text(slide, right_x + pad, col_top + Inches(0.85), right_w - pad * 2, Inches(0.3),
             "NVIDIA Corporation · Tech · $2.20T mkt cap",
             TextStyle(font=FONT_BODY, size_pt=11, color=TEXT_MUTED))
    # verdict pill
    add_rect(slide, right_x + pad, col_top + Inches(1.22), Inches(1.9), Inches(0.32),
             fill=RGBColor(0x12, 0x2E, 0x1B), border=POSITIVE, border_w_pt=0.5, radius=0.5)
    add_text(slide, right_x + pad, col_top + Inches(1.22), Inches(1.9), Inches(0.32),
             "✓  AAOIFI COMPLIANT",
             TextStyle(font=FONT_MONO, size_pt=10, bold=True, color=POSITIVE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE))

    # screen rows
    rows = [
        ("1. Debt / Market Cap",        "4.1%"),
        ("2. Interest Income / Revenue","0.0%"),
        ("3. Cash & AR / Market Cap",   "14.8%"),
        ("4. Haram industry flag",      "Clean"),
    ]
    row_y0 = col_top + Inches(1.85)
    each = Inches(0.34)
    for i, (lab, val) in enumerate(rows):
        y = row_y0 + each * i
        add_text(slide, right_x + pad, y, right_w / 2, Inches(0.3),
                 lab,
                 TextStyle(font=FONT_BODY, size_pt=11, color=TEXT_SECONDARY))
        add_rich_text(
            slide, right_x + right_w / 2, y, right_w / 2 - pad, Inches(0.3),
            runs=[
                (val + "  ", TextStyle(font=FONT_MONO, size_pt=13, bold=True, color=TEXT_PRIMARY)),
                ("✓ pass",   TextStyle(font=FONT_MONO, size_pt=11, color=POSITIVE)),
            ],
            align=PP_ALIGN.RIGHT,
        )
        # divider
        add_line(slide, right_x + pad, y + each - Inches(0.04), right_x + right_w - pad, y + each - Inches(0.04),
                 color=BORDER_SUBTLE, width_pt=0.4)

    add_text(slide, right_x + pad, col_top + spec_h - Inches(0.4), right_w - pad * 2, Inches(0.3),
             "Source · FMP Fundamentals 2026-Q1 · verified by halal_screening.py",
             TextStyle(font=FONT_MONO, size_pt=9, color=TEXT_MUTED))

    draw_slide_chrome(slide, 3)


def build_backtest_slide(prs: Presentation) -> None:
    """Slide 04 — strategy backtest comparison (A · B winner · C)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)

    # Header
    draw_eyebrow(slide, PAD_X, Inches(0.55), "Strategy backtest · comparison")
    add_rich_text(
        slide, PAD_X, Inches(0.95), Inches(12), Inches(1.2),
        runs=[
            ("Three strategies, one halal universe — ",
             TextStyle(font=FONT_DISPLAY, size_pt=32, bold=True, color=TEXT_PRIMARY, line_spacing=1.08)),
            ("five-year walk-forward.",
             TextStyle(font=FONT_DISPLAY, size_pt=32, bold=True, color=ACCENT, line_spacing=1.08)),
        ],
        line_spacing=1.08,
    )
    add_text(
        slide, PAD_X, Inches(1.95), Inches(12), Inches(0.4),
        "2021 – 2026 · S&P 500 post-AAOIFI screening · 20bps round-trip · half-Kelly sizing.",
        TextStyle(font=FONT_DISPLAY, size_pt=13, color=TEXT_SECONDARY),
    )

    # 3 columns
    strats = [
        ("Strategy A", "strat_A_breakout  ·  142 trades", "+18.4%", "1.62", "−9.4%",  "58%", "6.5%", False),
        ("Strategy B", "strat_B_mean_rev  ·   98 trades", "+22.8%", "1.84", "−6.1%",  "64%", "5.0%", True),
        ("Strategy C", "strat_C_momentum ·   64 trades",  "+15.2%", "1.21", "−11.8%", "51%", "6.3%", False),
    ]
    grid_top = Inches(2.6)
    grid_h   = Inches(4.0)
    gap = Inches(0.22)
    col_w = (SLIDE_W - PAD_X * 2 - gap * 2) / 3

    for i, (name, desc, cagr, sharpe, dd, win, kelly, winner) in enumerate(strats):
        x = PAD_X + (col_w + gap) * i
        # card
        border_color = ACCENT if winner else BORDER_DEFAULT
        card = add_rect(slide, x, grid_top, col_w, grid_h, fill=BG_SURFACE, border=border_color,
                        border_w_pt=1.0 if winner else 0.5, radius=0.04)

        # "Selected" tag on winner
        if winner:
            tag_w = Inches(1.0)
            add_rect(slide, x + col_w - tag_w - Inches(0.3), grid_top, tag_w, Inches(0.28),
                     fill=ACCENT, border=None, radius=0)
            add_text(slide, x + col_w - tag_w - Inches(0.3), grid_top - Inches(0.02), tag_w, Inches(0.28),
                     "SELECTED",
                     TextStyle(font=FONT_MONO, size_pt=8, bold=True, color=BG_ROOT,
                               align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE))

        pad = Inches(0.32)
        add_text(slide, x + pad, grid_top + Inches(0.3), col_w - pad * 2, Inches(0.5),
                 name,
                 TextStyle(font=FONT_DISPLAY, size_pt=18, bold=True, color=TEXT_PRIMARY))
        add_text(slide, x + pad, grid_top + Inches(0.78), col_w - pad * 2, Inches(0.3),
                 desc,
                 TextStyle(font=FONT_MONO, size_pt=10, color=TEXT_MUTED))

        # Equity curve placeholder (small chart card)
        chart_y = grid_top + Inches(1.2)
        chart_h = Inches(1.0)
        add_rect(slide, x + pad, chart_y, col_w - pad * 2, chart_h,
                 fill=BG_RAISED, border=BORDER_SUBTLE, border_w_pt=0.5, radius=0.03)
        # Sketch a curve as a connected polyline using line shapes
        points_a = [(0.05, 0.85), (0.18, 0.78), (0.32, 0.70), (0.46, 0.60), (0.60, 0.48),
                    (0.74, 0.36), (0.86, 0.25), (0.95, 0.18)]
        points_b = [(0.05, 0.88), (0.18, 0.74), (0.32, 0.60), (0.46, 0.45), (0.60, 0.34),
                    (0.74, 0.22), (0.86, 0.13), (0.95, 0.07)]
        points_c = [(0.05, 0.82), (0.18, 0.75), (0.32, 0.65), (0.46, 0.70), (0.60, 0.55),
                    (0.74, 0.45), (0.86, 0.36), (0.95, 0.30)]
        points = [points_a, points_b, points_c][i]
        cx0 = x + pad + Inches(0.05)
        cy0 = chart_y + Inches(0.05)
        cw  = col_w - pad * 2 - Inches(0.1)
        ch  = chart_h - Inches(0.1)
        for j in range(len(points) - 1):
            x1 = cx0 + cw * points[j][0]
            y1 = cy0 + ch * points[j][1]
            x2 = cx0 + cw * points[j + 1][0]
            y2 = cy0 + ch * points[j + 1][1]
            add_line(slide, x1, y1, x2, y2, color=ACCENT, width_pt=2.0 if winner else 1.6)

        # metrics 2-up
        metrics_y = grid_top + Inches(2.4)
        mh = Inches(0.55)
        mw = (col_w - pad * 2 - Inches(0.15)) / 2

        def metric_tile(lx, ly, lab, val, color=TEXT_PRIMARY):
            add_text(slide, lx, ly, mw, Inches(0.22), lab.upper(),
                     TextStyle(font=FONT_BODY, size_pt=8, bold=True, color=TEXT_MUTED))
            add_text(slide, lx, ly + Inches(0.22), mw, Inches(0.4), val,
                     TextStyle(font=FONT_MONO, size_pt=22, bold=True, color=color, line_spacing=1.0))

        metric_tile(x + pad,                 metrics_y,            "CAGR",  cagr, POSITIVE)
        metric_tile(x + pad + mw + Inches(0.15), metrics_y,        "Sharpe", sharpe)
        metric_tile(x + pad,                 metrics_y + mh + Inches(0.2),  "Max DD", dd, NEGATIVE)
        metric_tile(x + pad + mw + Inches(0.15), metrics_y + mh + Inches(0.2), "Win rate", win)
        # full-width kelly row
        add_line(slide, x + pad, metrics_y + mh * 2 + Inches(0.45),
                       x + col_w - pad, metrics_y + mh * 2 + Inches(0.45),
                       color=BORDER_SUBTLE, width_pt=0.4)
        add_text(slide, x + pad, metrics_y + mh * 2 + Inches(0.5), mw, Inches(0.22),
                 "KELLY f USED",
                 TextStyle(font=FONT_BODY, size_pt=8, bold=True, color=TEXT_MUTED))
        add_text(slide, x + col_w - pad - Inches(2), metrics_y + mh * 2 + Inches(0.5), Inches(2), Inches(0.3),
                 kelly,
                 TextStyle(font=FONT_MONO, size_pt=16, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT))

    # Footnote
    add_text(
        slide,
        PAD_X, SLIDE_W - Inches(7.85) if False else Inches(6.85),
        SLIDE_W - PAD_X * 2, Inches(0.3),
        "Past performance does not guarantee future results. Backtests assume 20bps round-trip cost, half-Kelly sizing, no shorting. Out-of-sample 2025-2026 matches in-sample within 0.15 Sharpe.",
        TextStyle(font=FONT_MONO, size_pt=9, color=TEXT_MUTED, align=PP_ALIGN.CENTER),
    )

    draw_slide_chrome(slide, 4)


def build_risk_slide(prs: Presentation) -> None:
    """Slide 05 — Risk Controls (Kelly narrative + 4 metric tiles)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)

    draw_eyebrow(slide, PAD_X, Inches(0.55), "Risk controls")
    add_rich_text(
        slide, PAD_X, Inches(0.95), Inches(11.5), Inches(1.4),
        runs=[
            ("Risk is ",                    TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=TEXT_PRIMARY, line_spacing=1.08)),
            ("mathematical",                TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=ACCENT,        line_spacing=1.08)),
            (", not editorial.",            TextStyle(font=FONT_DISPLAY, size_pt=36, bold=True, color=TEXT_PRIMARY, line_spacing=1.08)),
        ],
        line_spacing=1.08,
    )

    # Two columns
    col_top = Inches(2.5)
    col_h   = Inches(4.0)
    narrow_w = Inches(5.2)
    wide_x = PAD_X + narrow_w + Inches(0.4)
    wide_w = SLIDE_W - PAD_X - wide_x

    # LEFT — narrative card
    add_rect(slide, PAD_X, col_top, narrow_w, col_h, fill=BG_SURFACE, border=BORDER_DEFAULT, border_w_pt=0.5, radius=0.04)
    npad = Inches(0.4)
    # icon tile
    add_rect(slide, PAD_X + npad, col_top + npad, Inches(0.6), Inches(0.6),
             fill=ACCENT_DIM, border=ACCENT, border_w_pt=0.5, radius=0.1)
    add_text(slide, PAD_X + npad, col_top + npad, Inches(0.6), Inches(0.6),
             "ƒ",
             TextStyle(font=FONT_DISPLAY, size_pt=24, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE))
    add_text(slide, PAD_X + npad, col_top + npad + Inches(0.85), narrow_w - npad * 2, Inches(0.5),
             "Continuous Kelly · position sizing",
             TextStyle(font=FONT_DISPLAY, size_pt=20, bold=True, color=TEXT_PRIMARY))
    add_text(slide, PAD_X + npad, col_top + npad + Inches(1.4), narrow_w - npad * 2, Inches(0.9),
             "For a strategy with per-trade returns r, the Kelly-optimal fraction maximizes long-run growth. We use half-Kelly by default — roughly 75% of full-Kelly growth with ~50% of the drawdown.",
             TextStyle(font=FONT_BODY, size_pt=12, color=TEXT_SECONDARY, line_spacing=1.45))
    # formula box
    add_rect(slide, PAD_X + npad, col_top + npad + Inches(2.4), Inches(3.5), Inches(0.5),
             fill=BG_RAISED, border=BORDER_SUBTLE, border_w_pt=0.5, radius=0.1)
    add_text(slide, PAD_X + npad, col_top + npad + Inches(2.4), Inches(3.5), Inches(0.5),
             "f* = mean(r) / var(r)",
             TextStyle(font=FONT_MONO, size_pt=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE))
    add_text(slide, PAD_X + npad, col_top + npad + Inches(3.0), narrow_w - npad * 2, Inches(0.7),
             "Sample-size shrinkage applies when n < 30. Static risk budget overrides Kelly when more conservative. A hard 25% ceiling is never crossed.",
             TextStyle(font=FONT_BODY, size_pt=10, color=TEXT_MUTED, line_spacing=1.45))

    # RIGHT — 4 metric tiles (2x2)
    tiles = [
        ("Portfolio VaR · 95%", "−2.5%", "One-day · live book · 30-sec tick", TEXT_PRIMARY),
        ("CVaR · tail",         "−3.8%", "Expected loss in worst 5%",          TEXT_PRIMARY),
        ("Sharpe · 30d",        "1.84",  "Across A · B · C book",              POSITIVE),
        ("Max DD",              "−8.2%", "Kill-switch arms at −15%",           TEXT_PRIMARY),
    ]
    tile_gap = Inches(0.18)
    tile_w = (wide_w - tile_gap) / 2
    tile_h = (col_h - tile_gap) / 2
    for i, (lab, val, sub, color) in enumerate(tiles):
        row, col = divmod(i, 2)
        tx = wide_x + col * (tile_w + tile_gap)
        ty = col_top + row * (tile_h + tile_gap)
        add_rect(slide, tx, ty, tile_w, tile_h, fill=BG_SURFACE, border=BORDER_DEFAULT, border_w_pt=0.5, radius=0.04)
        tp = Inches(0.32)
        add_text(slide, tx + tp, ty + tp, tile_w - tp * 2, Inches(0.3),
                 lab.upper(),
                 TextStyle(font=FONT_BODY, size_pt=10, bold=True, color=TEXT_MUTED))
        add_text(slide, tx + tp, ty + Inches(0.7), tile_w - tp * 2, Inches(1.0),
                 val,
                 TextStyle(font=FONT_MONO, size_pt=44, bold=True, color=color, line_spacing=1.0))
        add_text(slide, tx + tp, ty + tile_h - Inches(0.5), tile_w - tp * 2, Inches(0.3),
                 sub,
                 TextStyle(font=FONT_BODY, size_pt=10, color=TEXT_SECONDARY))

    draw_slide_chrome(slide, 5)


def build_closing_slide(prs: Presentation) -> None:
    """Slide 06 — Closing CTA."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_slide_bg(slide)

    # Subtle glow band at bottom (approximation via large dim rectangle)
    add_rect(slide, Inches(0), SLIDE_H - Inches(4), SLIDE_W, Inches(4),
             fill=RGBColor(0x16, 0x12, 0x14), border=None, radius=0)

    # Centered brand mark
    mark_size = Inches(0.95)
    cx = (SLIDE_W - mark_size) / 2
    draw_brand_mark(slide, cx, Inches(1.4), size=mark_size)

    # Headline
    add_rich_text(
        slide, Inches(0), Inches(2.8), SLIDE_W, Inches(2.0),
        runs=[
            ("A small number of\n", TextStyle(font=FONT_DISPLAY, size_pt=52, bold=True, color=TEXT_PRIMARY, line_spacing=1.05, align=PP_ALIGN.CENTER)),
            ("institutional partners",   TextStyle(font=FONT_DISPLAY, size_pt=52, bold=True, color=ACCENT,        line_spacing=1.05, align=PP_ALIGN.CENTER)),
            (" each quarter.",           TextStyle(font=FONT_DISPLAY, size_pt=52, bold=True, color=TEXT_PRIMARY, line_spacing=1.05, align=PP_ALIGN.CENTER)),
        ],
        align=PP_ALIGN.CENTER,
        line_spacing=1.05,
    )

    add_text(slide, Inches(0), Inches(4.95), SLIDE_W, Inches(0.5),
             "عدد محدود من الشركاء المؤسسيين كل ربع سنة.",
             TextStyle(font=FONT_ARABIC, size_pt=20, bold=True, color=TEXT_SECONDARY, align=PP_ALIGN.CENTER))

    # 3-block contact strip
    blocks = [
        ("BRIEFING", "60 minutes · live demo + Q&A"),
        ("LANGUAGES", "English · العربية"),
        ("CONTACT",   "institutional@mizanquant.com"),
    ]
    strip_y = Inches(5.9)
    bw = Inches(3.4)
    gap = Inches(0.6)
    total = bw * 3 + gap * 2
    start_x = (SLIDE_W - total) / 2
    for i, (label, value) in enumerate(blocks):
        x = start_x + (bw + gap) * i
        add_text(slide, x, strip_y, bw, Inches(0.25),
                 label,
                 TextStyle(font=FONT_BODY, size_pt=9, bold=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER))
        add_text(slide, x, strip_y + Inches(0.28), bw, Inches(0.4),
                 value,
                 TextStyle(font=FONT_DISPLAY, size_pt=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER))


# =============================================================================
# MAIN
# =============================================================================

def build_deck(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs)
    build_architecture_slide(prs)
    build_halal_slide(prs)
    build_backtest_slide(prs)
    build_risk_slide(prs)
    build_closing_slide(prs)

    prs.save(str(out_path))
    return out_path


def main() -> None:
    p = argparse.ArgumentParser(description="Export the MizanQuant investor briefing as a native .pptx")
    p.add_argument(
        "--out",
        default=str(Path(__file__).resolve().parent / "mizanquant-briefing.pptx"),
        help="Output .pptx path (default: alongside this script)",
    )
    args = p.parse_args()
    out = build_deck(Path(args.out))
    print(f"✓ Wrote {out}")
    print("  6 slides · native editable shapes · 16:9 · DM Sans + JetBrains Mono + Cairo")
    print("  Open in PowerPoint → File → Options → Save → Embed fonts (Cairo / DM Sans)")


if __name__ == "__main__":
    main()
